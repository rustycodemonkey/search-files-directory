from pptx import Presentation
import glob

# for eachfile in glob.glob("/Users/phillip.kim/Box/CV Blitz/[J-Z]*.pptx"):
for eachfile in glob.glob("/Users/phillip.kim/Library/CloudStorage/Box-Box/CV Bios/*.pptx"):
    try:
        prs = Presentation(eachfile)
        print(eachfile)
        print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except:
        print(eachfile)
        print("ERROR")
        print("----------------------")
